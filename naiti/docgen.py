"""
docgen.py — render generated content into real Office/text files.

Every file naiti uploads is a genuine .docx / .xlsx / .pptx / .csv / .txt, so
the upload path exercises the app's real extractors rather than a shortcut.
Files are built in memory and posted straight to the API; nothing is written
to disk.
"""
from __future__ import annotations

import csv
import io


def txt_bytes(body: str) -> bytes:
    return body.encode("utf-8")


def csv_bytes(headers: list[str], rows: list[list]) -> bytes:
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(headers)
    for r in rows:
        w.writerow(r)
    return buf.getvalue().encode("utf-8")


def docx_bytes(title: str, blocks: list[dict]) -> bytes:
    """blocks: [{'heading': str}, {'para': str}, {'bullets': [str]},
                {'table': {'headers': [...], 'rows': [[...]]}}]"""
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)

    for b in blocks:
        if "heading" in b:
            doc.add_heading(b["heading"], level=b.get("level", 1))
        if "para" in b:
            doc.add_paragraph(b["para"])
        if "bullets" in b:
            for item in b["bullets"]:
                doc.add_paragraph(item, style="List Bullet")
        if "table" in b:
            t = b["table"]
            table = doc.add_table(rows=1, cols=len(t["headers"]))
            table.style = "Light Grid Accent 1"
            for i, h in enumerate(t["headers"]):
                table.rows[0].cells[i].text = str(h)
            for row in t["rows"]:
                cells = table.add_row().cells
                for i, v in enumerate(row):
                    cells[i].text = str(v)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def xlsx_bytes(sheets: list[dict]) -> bytes:
    """sheets: [{'name': str, 'headers': [...], 'rows': [[...]], 'title': str?}]"""
    import openpyxl
    from openpyxl.styles import Font

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    for s in sheets:
        ws = wb.create_sheet(title=s["name"][:31])
        r = 1
        if s.get("title"):
            ws.cell(row=1, column=1, value=s["title"]).font = Font(bold=True, size=13)
            r = 3
        for i, h in enumerate(s["headers"], start=1):
            ws.cell(row=r, column=i, value=h).font = Font(bold=True)
        for row in s["rows"]:
            r += 1
            for i, v in enumerate(row, start=1):
                ws.cell(row=r, column=i, value=v)
        for i, h in enumerate(s["headers"], start=1):
            ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = max(14, len(str(h)) + 3)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def pptx_bytes(title: str, subtitle: str, slides: list[dict]) -> bytes:
    """slides: [{'title': str, 'bullets': [str]}]"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title
    cover.placeholders[1].text = subtitle

    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s["title"]
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(s.get("bullets", [])):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = line
            p.font.size = Pt(16)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def money(v: float) -> str:
    """Format like a finance team would: EUR 18,412,000."""
    return f"EUR {round(v):,}"


def millions(v: float) -> str:
    return f"EUR {v / 1_000_000:.1f}M"


def pct(v: float) -> str:
    return f"{v:.1f}%"
